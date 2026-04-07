from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt
from playwright.sync_api import sync_playwright


BASE_URL = "http://127.0.0.1:3000"
ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "20260403"
ASSET_DIR = OUT_DIR / "ppt_assets"
PPT_PATH = OUT_DIR / "20260403-系统演示与PPT.pptx"
FLOW_PATH = OUT_DIR / "20260403-演示流程.md"

USERNAME = "admin"
PASSWORD = "admin123"
VIEWPORT = {"width": 1600, "height": 900}

NAVY = RGBColor(15, 50, 76)
BLUE = RGBColor(47, 140, 173)
BLUE_DARK = RGBColor(31, 111, 151)
LIGHT_BG = RGBColor(245, 249, 252)
CARD_BG = RGBColor(255, 255, 255)
TEXT = RGBColor(28, 41, 53)
MUTED = RGBColor(98, 110, 128)
GREEN = RGBColor(45, 130, 81)
ORANGE = RGBColor(190, 110, 30)
RED = RGBColor(186, 58, 60)
LINE = RGBColor(219, 232, 242)


@dataclass
class Shot:
    title: str
    path: Path
    caption: str


def ensure_dirs() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def get_captcha(page) -> str:
    return page.locator("span.login-captcha-box").inner_text().strip()


def login(page) -> None:
    page.goto(f"{BASE_URL}/login", wait_until="networkidle")
    captcha = get_captcha(page)
    page.fill('input[name="username"]', USERNAME)
    page.fill('input[name="password"]', PASSWORD)
    page.fill('input[name="captcha_input"]', captcha)
    page.get_by_role("button", name="登录").click()
    page.wait_for_url("**/workbench")
    page.wait_for_load_state("networkidle")


def first_href(page, selector: str) -> str:
    loc = page.locator(selector).first
    href = loc.get_attribute("href")
    if not href:
        raise RuntimeError(f"Missing href for selector: {selector}")
    return href


def wait_ready(page, selector: str | None = None) -> None:
    page.wait_for_load_state("networkidle")
    if selector:
        page.locator(selector).wait_for(state="visible", timeout=15000)


def capture(page, url: str, out_path: Path, selector: str | None = None, settle_ms: int = 1200) -> None:
    page.goto(url, wait_until="networkidle")
    if selector:
        page.locator(selector).wait_for(state="visible", timeout=15000)
    page.wait_for_timeout(settle_ms)
    page.screenshot(path=str(out_path), full_page=False)


def capture_all() -> list[Shot]:
    shots: list[Shot] = []
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        context = browser.new_context(viewport=VIEWPORT, device_scale_factor=1)
        page = context.new_page()

        login(page)
        login_shot = ASSET_DIR / "01-login.png"
        page.goto(f"{BASE_URL}/login", wait_until="networkidle")
        page.screenshot(path=str(login_shot), full_page=False)
        shots.append(Shot("登录页", login_shot, "开场：展示系统入口与品牌识别。"))

        dashboard_shot = ASSET_DIR / "02-dashboard.png"
        capture(page, f"{BASE_URL}/", dashboard_shot, ".crm-page-head")
        shots.append(Shot("总览", dashboard_shot, "展示经营概览、异常看板和项目状态。"))

        workbench_shot = ASSET_DIR / "03-workbench.png"
        capture(page, f"{BASE_URL}/workbench", workbench_shot, ".crm-workbench-hero", 1800)
        shots.append(Shot("工作台", workbench_shot, "展示待办、关注事项与 AI 优先建议。"))

        customers_shot = ASSET_DIR / "04-customers.png"
        capture(page, f"{BASE_URL}/customers", customers_shot, ".crm-data-card")
        shots.append(Shot("客户列表", customers_shot, "客户经营入口：筛选、导出和新增。"))

        opportunities_shot = ASSET_DIR / "05-opportunities.png"
        capture(page, f"{BASE_URL}/opportunities", opportunities_shot, ".crm-data-card")
        shots.append(Shot("商机列表", opportunities_shot, "销售漏斗与跟进节奏入口。"))

        contracts_shot = ASSET_DIR / "06-contracts.png"
        capture(page, f"{BASE_URL}/contracts", contracts_shot, ".crm-data-card")
        shots.append(Shot("合同列表", contracts_shot, "签约与履约的核心台账。"))

        projects_list_shot = ASSET_DIR / "07-projects.png"
        capture(page, f"{BASE_URL}/projects", projects_list_shot, ".crm-data-card")
        shots.append(Shot("项目列表", projects_list_shot, "从合同落到项目交付的枢纽。"))

        project_href = first_href(page, 'table tbody a[href^="/projects/"]')
        project_detail_shot = ASSET_DIR / "08-project-detail.png"
        capture(page, f"{BASE_URL}{project_href}", project_detail_shot, ".project-detail-page", 1600)
        shots.append(Shot("项目详情", project_detail_shot, "进度、任务、里程碑、风险与结项审批。"))

        receivables_shot = ASSET_DIR / "09-receivables.png"
        capture(page, f"{BASE_URL}/receivables", receivables_shot, ".crm-data-card")
        shots.append(Shot("回款应收", receivables_shot, "项目、合同与财务协同的现金流视图。"))

        invoices_shot = ASSET_DIR / "10-invoices.png"
        capture(page, f"{BASE_URL}/invoices", invoices_shot, ".crm-data-card")
        shots.append(Shot("开票管理", invoices_shot, "回款之后的开票归档与管理。"))

        approvals_shot = ASSET_DIR / "11-approvals.png"
        capture(page, f"{BASE_URL}/approvals", approvals_shot, ".crm-data-card")
        shots.append(Shot("审批中心", approvals_shot, "赢单、签约、结项审批的统一入口。"))

        browser.close()
    return shots


def add_box(slide, left, top, width, height, fill=CARD_BG, line=LINE, radius=0.16):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, bullets, font_size=16, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        p.text = bullet
        for run in p.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(10.5), Inches(0.45), title, font_size=24, color=NAVY, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.68), Inches(11), Inches(0.28), subtitle, font_size=11, color=MUTED)


def add_screenshot_card(slide, shot: Shot, left, top, width, height, caption_height=0.38):
    add_box(slide, left, top, width, height, fill=CARD_BG, line=LINE)
    img_margin = Inches(0.08)
    pic_top = top + img_margin
    pic_left = left + img_margin
    pic_width = width - img_margin * 2
    pic_height = height - img_margin * 2 - Inches(caption_height)
    slide.shapes.add_picture(str(shot.path), pic_left, pic_top, width=pic_width, height=pic_height)
    add_textbox(
        slide,
        left + Inches(0.15),
        top + height - Inches(caption_height) - Inches(0.03),
        width - Inches(0.3),
        Inches(caption_height),
        shot.title + " · " + shot.caption,
        font_size=11,
        color=MUTED,
    )


def add_top_tag(slide, text):
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.95), Inches(2.7), Inches(0.38))
    tag.fill.solid()
    tag.fill.fore_color.rgb = RGBColor(230, 243, 249)
    tag.line.color.rgb = RGBColor(195, 226, 238)
    add_textbox(slide, Inches(0.62), Inches(1.0), Inches(2.45), Inches(0.2), text, font_size=11, color=BLUE_DARK, bold=True)


def build_cover(prs, shot: Shot):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_box(slide, Inches(0.45), Inches(0.4), Inches(12.4), Inches(6.6), fill=RGBColor(248, 251, 253), line=LINE)
    add_top_tag(slide, "2026-04-03 下午演示")
    add_textbox(slide, Inches(0.6), Inches(1.45), Inches(6.2), Inches(1.0), "项目过程管理系统", font_size=28, color=NAVY, bold=True)
    add_textbox(slide, Inches(0.6), Inches(2.35), Inches(6.0), Inches(0.6), "全局功能主讲，权限矩阵和审批流程为辅", font_size=18, color=BLUE_DARK, bold=True)
    add_bullets(
        slide,
        Inches(0.62),
        Inches(3.05),
        Inches(5.6),
        Inches(1.8),
        [
            "一个入口看经营：工作台、总览、客户、商机、合同、项目、回款、开票。",
            "一个链路讲闭环：从销售推进到项目交付，再回到财务收口与审批流转。",
            "一个收尾讲治理：8 角色权限边界 + 3 类审批事项，说明系统可控可管。",
        ],
        font_size=15,
        color=TEXT,
    )
    add_screenshot_card(slide, shot, Inches(7.0), Inches(1.2), Inches(5.3), Inches(4.95), caption_height=0.42)
    add_textbox(slide, Inches(7.2), Inches(6.35), Inches(4.8), Inches(0.25), "登录页作为开场，先把系统品牌和入口感拉出来。", font_size=10, color=MUTED)


def build_screenshot_slide(prs, title, subtitle, shot: Shot, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_title(slide, title, subtitle)
    add_screenshot_card(slide, shot, Inches(0.5), Inches(1.15), Inches(8.2), Inches(5.95), caption_height=0.42)
    add_box(slide, Inches(8.9), Inches(1.15), Inches(3.95), Inches(5.95), fill=CARD_BG, line=LINE)
    add_textbox(slide, Inches(9.15), Inches(1.45), Inches(3.4), Inches(0.35), "现场讲解要点", font_size=18, color=NAVY, bold=True)
    add_bullets(slide, Inches(9.18), Inches(1.95), Inches(3.3), Inches(4.6), bullets, font_size=15, color=TEXT)


def build_dual_slide(prs, title, subtitle, shot_left: Shot, shot_right: Shot, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_title(slide, title, subtitle)
    add_screenshot_card(slide, shot_left, Inches(0.45), Inches(1.15), Inches(6.15), Inches(4.95), caption_height=0.38)
    add_screenshot_card(slide, shot_right, Inches(6.75), Inches(1.15), Inches(6.15), Inches(4.95), caption_height=0.38)
    add_box(slide, Inches(0.45), Inches(6.25), Inches(12.45), Inches(0.72), fill=RGBColor(236, 244, 249), line=RGBColor(213, 230, 241))
    add_bullets(slide, Inches(0.65), Inches(6.38), Inches(12.0), Inches(0.42), bullets, font_size=14, color=TEXT)


def build_project_slide(prs, shot: Shot):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_title(slide, "项目交付", "项目详情页最适合讲“管理动作如何落到执行层”。")
    add_screenshot_card(slide, shot, Inches(0.45), Inches(1.15), Inches(8.0), Inches(5.95), caption_height=0.42)
    add_box(slide, Inches(8.75), Inches(1.15), Inches(4.1), Inches(5.95), fill=CARD_BG, line=LINE)
    add_textbox(slide, Inches(9.0), Inches(1.4), Inches(3.4), Inches(0.35), "项目详情页可讲的重点", font_size=18, color=NAVY, bold=True)
    add_bullets(
        slide,
        Inches(9.0),
        Inches(1.9),
        Inches(3.45),
        Inches(4.9),
        [
            "项目概览把进度、阶段、风险和组织结构放在同一屏。",
            "任务、里程碑、风险和进展汇报构成真正的交付闭环。",
            "结项审批把项目经理、审批人和说明情况串成标准流程。",
            "这个页面最适合在演示里说明：销售成交后，交付不是结束，而是开始。",
        ],
        font_size=15,
        color=TEXT,
    )


def build_matrix_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_title(slide, "权限矩阵与审批流程", "作为辅助内容，用一页讲清“谁能看、谁能改、谁能批”。")

    role_cards = [
        ("admin", "最高权限，系统维护和全局兜底"),
        ("management", "只看全局经营，不做业务处理"),
        ("sales", "经营客户、商机、合同与回款"),
        ("pm", "管理自己负责的项目交付过程"),
        ("implementer", "仅执行本人参与项目中的任务与记录"),
        ("finance", "全局只读，重点管理开票"),
        ("sales_director", "销售线管理与赢单/签约审批"),
        ("project_director", "项目线管理与结项审批"),
    ]

    start_x = Inches(0.45)
    start_y = Inches(1.15)
    card_w = Inches(3.0)
    card_h = Inches(0.9)
    gap_x = Inches(0.18)
    gap_y = Inches(0.16)
    for idx, (role, desc) in enumerate(role_cards):
        col = idx % 4
        row = idx // 4
        left = start_x + col * (card_w + gap_x)
        top = start_y + row * (card_h + gap_y)
        add_box(slide, left, top, card_w, card_h, fill=CARD_BG, line=LINE)
        add_textbox(slide, left + Inches(0.12), top + Inches(0.08), Inches(0.9), Inches(0.18), role, font_size=12, color=BLUE_DARK, bold=True)
        add_textbox(slide, left + Inches(0.12), top + Inches(0.3), Inches(2.68), Inches(0.38), desc, font_size=10.5, color=TEXT)

    add_box(slide, Inches(0.45), Inches(3.35), Inches(12.45), Inches(2.45), fill=CARD_BG, line=LINE)
    add_textbox(slide, Inches(0.65), Inches(3.58), Inches(3.0), Inches(0.28), "审批流程", font_size=18, color=NAVY, bold=True)

    flow = [
        ("商机赢单", "销售发起", "销售总监审批"),
        ("合同签约", "销售发起", "销售总监审批"),
        ("项目结项", "项目经理发起", "项目总监审批"),
    ]
    y = 4.0
    for left_label, start_label, end_label in flow:
        add_box(slide, Inches(0.7), Inches(y), Inches(3.0), Inches(0.48), fill=RGBColor(234, 243, 249), line=RGBColor(200, 226, 238))
        add_textbox(slide, Inches(0.85), Inches(y + 0.11), Inches(1.0), Inches(0.18), left_label, font_size=12, color=BLUE_DARK, bold=True)
        add_textbox(slide, Inches(2.05), Inches(y + 0.11), Inches(0.75), Inches(0.18), start_label, font_size=11, color=TEXT)
        add_textbox(slide, Inches(3.0), Inches(y + 0.08), Inches(0.3), Inches(0.18), "→", font_size=14, color=BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_box(slide, Inches(3.35), Inches(y), Inches(2.4), Inches(0.48), fill=RGBColor(243, 248, 252), line=LINE)
        add_textbox(slide, Inches(3.5), Inches(y + 0.11), Inches(2.0), Inches(0.18), end_label, font_size=11, color=TEXT)
        y += 0.64

    add_textbox(
        slide,
        Inches(6.2),
        Inches(3.62),
        Inches(6.2),
        Inches(1.55),
        "三层权限是这套系统最容易讲清的治理点：\n1. 菜单权限决定看不看得到\n2. 操作权限决定能不能做\n3. 数据权限决定看多大范围\n\n现场讲这页时，重点强调“业务链路清晰，但权限边界收得住”。",
        font_size=13,
        color=TEXT,
    )


def build_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_title(slide, "收尾总结", "最后用三句话把演示价值收回来。")
    add_box(slide, Inches(0.55), Inches(1.25), Inches(3.95), Inches(4.9), fill=CARD_BG, line=LINE)
    add_box(slide, Inches(4.72), Inches(1.25), Inches(3.95), Inches(4.9), fill=CARD_BG, line=LINE)
    add_box(slide, Inches(8.9), Inches(1.25), Inches(3.95), Inches(4.9), fill=CARD_BG, line=LINE)

    add_textbox(slide, Inches(0.82), Inches(1.6), Inches(3.4), Inches(0.3), "1. 全局功能一张图", font_size=18, color=NAVY, bold=True)
    add_bullets(slide, Inches(0.82), Inches(2.05), Inches(3.2), Inches(3.6), [
        "总览、工作台和业务列表构成一条完整的入口链。",
        "演示时先讲整体，再切具体模块，节奏最顺。",
    ], font_size=15)

    add_textbox(slide, Inches(4.99), Inches(1.6), Inches(3.4), Inches(0.3), "2. 业务链路可闭环", font_size=18, color=NAVY, bold=True)
    add_bullets(slide, Inches(4.99), Inches(2.05), Inches(3.2), Inches(3.6), [
        "客户、商机、合同、项目、回款、开票彼此串联。",
        "交付和财务不再是两个孤立系统，而是一条主线。",
    ], font_size=15)

    add_textbox(slide, Inches(9.17), Inches(1.6), Inches(3.4), Inches(0.3), "3. 权限与审批可控", font_size=18, color=NAVY, bold=True)
    add_bullets(slide, Inches(9.17), Inches(2.05), Inches(3.2), Inches(3.6), [
        "8 角色边界明确，避免“谁都能改”。",
        "赢单、签约、结项三条审批把关键动作收口。",
        "这部分是管理层最容易接受的安全感。",
    ], font_size=15)

    add_textbox(slide, Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.35), "建议现场讲法：先总览、再链路、后治理；每一页都围绕“能不能把业务跑起来”来讲。", font_size=12, color=BLUE_DARK, align=PP_ALIGN.CENTER)


def write_flow_doc(shots: list[Shot]) -> None:
    content = """# 2026-04-03 演示流程

## 讲解顺序
1. 登录页：先给出系统入口和品牌第一印象。
2. 总览页：说明系统能从经营、项目和财务三个维度看全局。
3. 工作台：展示待办、关注事项和 AI 优先建议，强调“今天要做什么”。
4. 客户与商机：说明销售如何推进客户经营与商机转化。
5. 合同与项目：说明签约之后如何自然落到项目交付。
6. 项目详情：展示任务、里程碑、风险、进展和结项审批。
7. 回款与开票：展示财务闭环如何和项目、合同联动。
8. 审批中心：说明赢单、签约、结项三类审批的统一入口。
9. 权限矩阵：用 8 个角色边界收尾，说明系统可控可管。
10. 总结页：回到系统价值和下一步落地建议。

## 现场表达重点
- 先讲“全局功能”，让观众先知道系统能解决什么问题。
- 再讲“业务链路”，突出客户到项目、项目到财务的闭环。
- 最后讲“权限和审批”，让管理层确认边界和风控方式。

## 已生成截图
"""
    for idx, shot in enumerate(shots, start=1):
        content += f"- {idx:02d}. {shot.title}：`ppt_assets/{shot.path.name}`\n"
    FLOW_PATH.write_text(content, encoding="utf-8")


def build_ppt(shots: list[Shot]) -> None:
    shot_map = {shot.title: shot for shot in shots}
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "2026-04-03 系统演示与PPT"
    prs.core_properties.subject = "项目过程管理系统"
    prs.core_properties.author = "Cursor AI"

    build_cover(prs, shot_map["登录页"])
    build_screenshot_slide(
        prs,
        "全局总览",
        "先让观众看懂系统整体运行状态，再进入业务细节。",
        shot_map["总览"],
        [
            "这里是全局经营入口，适合管理层先看大盘。",
            "异常看板把项目暂停、待审批和逾期应收一起拉出来。",
            "这一页的作用是告诉大家：系统不是单点功能，而是一个统一视图。",
        ],
    )
    build_screenshot_slide(
        prs,
        "工作台与节奏",
        "工作台更像每天开工前的指挥台。",
        shot_map["工作台"],
        [
            "待办总量、需关注事项和经营指标让节奏一眼可见。",
            "AI 优先建议是辅助，不是噱头，重点是帮人排序。",
            "这页最适合对销售、项目和管理层一起讲。",
        ],
    )
    build_dual_slide(
        prs,
        "客户与商机",
        "先经营客户，再推进商机，是销售链路最核心的前半段。",
        shot_map["客户列表"],
        shot_map["商机列表"],
        [
            "客户列表负责沉淀客户台账和筛选导出。",
            "商机列表承接转化过程，重点看阶段、金额和负责人。",
            "这两页一起讲，观众会很快理解“线索到成交”的路径。",
        ],
    )
    build_dual_slide(
        prs,
        "合同与项目",
        "签约之后就进入交付节奏，系统把前后链路接得很顺。",
        shot_map["合同列表"],
        shot_map["项目列表"],
        [
            "合同是收入和履约的关键节点。",
            "项目列表把签约结果转成交付管理对象。",
            "这一页适合强调销售到项目的交接不会断层。",
        ],
    )
    build_project_slide(prs, shot_map["项目详情"])
    build_dual_slide(
        prs,
        "回款与开票",
        "财务闭环是业务系统最容易被忽略、但最能体现完整性的部分。",
        shot_map["回款应收"],
        shot_map["开票管理"],
        [
            "回款应收把合同、项目和现金流串起来。",
            "开票管理负责记录发票号码、类型和状态。",
            "这页让观众看到：业务做完之后，账也能收得住。",
        ],
    )
    build_screenshot_slide(
        prs,
        "审批中心",
        "赢单、签约、结项三类审批统一收口。",
        shot_map["审批中心"],
        [
            "待办审批一眼可见，适合管理层和审批人快速处理。",
            "列表页已经把申请人、审批人、状态和时间线放在一起。",
            "如果需要强调治理能力，这页是最合适的收口页之一。",
        ],
    )
    build_matrix_slide(prs)
    build_closing_slide(prs)

    prs.save(PPT_PATH)


def main() -> int:
    ensure_dirs()
    shots = capture_all()
    write_flow_doc(shots)
    build_ppt(shots)
    print(f"Created: {PPT_PATH}")
    print(f"Created: {FLOW_PATH}")
    print(f"Screenshots: {ASSET_DIR}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
